"""
Kimi PPT 生成工具 —— 使用 PPTD 格式 + python-pptx 本地渲染导出 PPTX
PPTD 是一种基于 YAML 的幻灯片描述格式，AI 友好，支持复杂排版。
"""
import os
import re
import uuid
from pathlib import Path
from io import BytesIO

import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

from agent.tools import emit_progress

# 输出目录
ROOT = Path(__file__).resolve().parent.parent.parent.parent
OUTPUT_DIR = ROOT / "server" / "static" / "papers"
os.makedirs(OUTPUT_DIR, exist_ok=True)

# PPTD 画布标准: 960x540 (16:9)
PPTD_W = 960
PPTD_H = 540
# PPTX 标准宽屏: 13.333x7.5 inches
PPTX_W_INCHES = 13.333
PPTX_H_INCHES = 7.5

# 形状映射
SHAPE_MAP = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
    "roundrect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "roundedrectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "diamond": MSO_SHAPE.DIAMOND,
    "chevron": MSO_SHAPE.CHEVRON,
    "arrow": MSO_SHAPE.RIGHT_ARROW,
    "pentagon": MSO_SHAPE.PENTAGON,
    "hexagon": MSO_SHAPE.HEXAGON,
    "star": MSO_SHAPE.STAR_5_POINT,
    "star5": MSO_SHAPE.STAR_5_POINT,
    "star4": MSO_SHAPE.STAR_4_POINT,
    "star6": MSO_SHAPE.STAR_6_POINT,
    "star7": MSO_SHAPE.STAR_7_POINT,
    "star8": MSO_SHAPE.STAR_8_POINT,
    "parallelogram": MSO_SHAPE.PARALLELOGRAM,
    "trapezoid": MSO_SHAPE.TRAPEZOID,
    "plaque": MSO_SHAPE.PLAQUE,
    "pie": MSO_SHAPE.PIE,
    "arc": MSO_SHAPE.ARC,
    "line": MSO_SHAPE.LINE_INVERSE,
    "bentarrow": MSO_SHAPE.BENT_ARROW,
    "stripedrightarrow": MSO_SHAPE.STRIPED_RIGHT_ARROW,
    "notchedrightarrow": MSO_SHAPE.NOTCHED_RIGHT_ARROW,
    "blockarc": MSO_SHAPE.BLOCK_ARC,
    "donut": MSO_SHAPE.DONUT,
    "moon": MSO_SHAPE.MOON,
    "smiley": MSO_SHAPE.SMILEY_FACE,
    "heart": MSO_SHAPE.HEART,
    "lightning": MSO_SHAPE.LIGHTNING_BOLT,
    "sun": MSO_SHAPE.SUN,
    "cloud": MSO_SHAPE.CLOUD,
    "wave": MSO_SHAPE.WAVE,
    "doublewave": MSO_SHAPE.DOUBLE_WAVE,
    "plus": MSO_SHAPE.MATH_PLUS,
    "flowchartprocess": MSO_SHAPE.FLOWCHART_PROCESS,
    "flowchartdecision": MSO_SHAPE.FLOWCHART_DECISION,
    "flowchartdata": MSO_SHAPE.FLOWCHART_DATA,
    "flowchartpredefinedprocess": MSO_SHAPE.FLOWCHART_PREDEFINED_PROCESS,
    "flowchartinternalstorage": MSO_SHAPE.FLOWCHART_INTERNAL_STORAGE,
    "flowchartdocument": MSO_SHAPE.FLOWCHART_DOCUMENT,
    "flowchartmultidocument": MSO_SHAPE.FLOWCHART_MULTIDOCUMENT,
    "flowchartterminator": MSO_SHAPE.FLOWCHART_TERMINATOR,
    "flowchartpreparation": MSO_SHAPE.FLOWCHART_PREPARATION,
    "flowchartmanualinput": MSO_SHAPE.FLOWCHART_MANUAL_INPUT,
    "flowchartmanualoperation": MSO_SHAPE.FLOWCHART_MANUAL_OPERATION,
    "flowchartconnector": MSO_SHAPE.FLOWCHART_CONNECTOR,
    "flowchartcard": MSO_SHAPE.FLOWCHART_CARD,
    "flowchartsummingjunction": MSO_SHAPE.FLOWCHART_SUMMING_JUNCTION,
    "flowchartor": MSO_SHAPE.FLOWCHART_OR,
    "flowchartcollate": MSO_SHAPE.FLOWCHART_COLLATE,
    "flowchartsort": MSO_SHAPE.FLOWCHART_SORT,
    "flowchartextract": MSO_SHAPE.FLOWCHART_EXTRACT,
    "flowchartmerge": MSO_SHAPE.FLOWCHART_MERGE,
    "flowchartstoreddata": MSO_SHAPE.FLOWCHART_STORED_DATA,
    "flowchartsequentialaccess": MSO_SHAPE.FLOWCHART_SEQUENTIAL_ACCESS_STORAGE,
    "flowchartmagneticdisk": MSO_SHAPE.FLOWCHART_MAGNETIC_DISK,
    "flowchartdirectaccess": MSO_SHAPE.FLOWCHART_DIRECT_ACCESS_STORAGE,
    "flowchartdisplay": MSO_SHAPE.FLOWCHART_DISPLAY,
}

SCHEMA = {
    "type": "function",
    "tag": "PPT",
    "function": {
        "name": "generate_kimi_ppt",
        "description": (
            "使用 PPTD 格式生成专业的 PPT 演示文稿（.pptx），本地渲染无需网络。"
            "PPTD 是 YAML 格式的幻灯片描述语言，支持丰富的排版：文本、形状、图片、渐变等。"
            "\n\n"
            "## PPTD 格式说明\n"
            "画布标准：960x540（16:9），坐标原点在左上角。\n\n"
            "### 元素类型：\n"
            "**text** - 文本元素，支持 HTML 标签：<p>段落, <strong>粗体, <span style=\"...\">内联样式\n"
            "```yaml\n"
            "- type: text\n"
            "  position: [x, y]\n"
            "  size: [width, height]\n"
            "  content:\n"
            "    text: |\n"
            "      <p style=\"font-size:36px;font-weight:bold;color:#2B579A;\">标题</p>\n"
            "      <p style=\"font-size:18px;color:#666666;\">正文</p>\n"
            "```\n\n"
            "**shape** - 形状元素\n"
            "```yaml\n"
            "- type: shape\n"
            "  position: [x, y]\n"
            "  size: [width, height]\n"
            "  shape:\n"
            "    preset: rectangle  # 支持: rectangle, ellipse, roundRect, triangle, diamond, chevron, arrow, star, heart, cloud, line 等\n"
            "    fill: \"#2B579A\"\n"
            "    # 渐变填充:\n"
            "    # fill: [{color: \"#2B579A\", offset: 0}, {color: \"#E86A17\", offset: 1}]\n"
            "```\n\n"
            "**image** - 图片元素\n"
            "```yaml\n"
            "- type: image\n"
            "  position: [x, y]\n"
            "  size: [width, height]\n"
            "  image:\n"
            "    src: https://example.com/image.jpg\n"
            "```\n\n"
            "文本样式属性（<p> 和 <span> 的 style 中支持）：\n"
            "font-size: 像素字号 | font-weight: bold/normal | color: #RRGGBB\n"
            "text-align: left/center/right | font-family: 字体名\n"
            "line-height: 行高倍数 | background: 背景色\n"
            "font-style: italic | text-decoration: underline"
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "title": {
                    "type": "string",
                    "description": "PPT 的标题"
                },
                "theme_colors": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "主题配色，5个颜色：[主色, 强调色, 深色, 白色, 浅灰背景]。默认：['#2B579A', '#E86A17', '#333333', '#FFFFFF', '#F0F2F5']"
                },
                "slides": {
                    "type": "array",
                    "description": "幻灯片页面列表",
                    "items": {
                        "type": "object",
                        "properties": {
                            "filename": {
                                "type": "string",
                                "description": "页面文件名，如 cover.page, slide1.page"
                            },
                            "content": {
                                "type": "string",
                                "description": "该页的 PPTD YAML 内容（elements 数组）"
                            }
                        },
                        "required": ["filename", "content"]
                    }
                }
            },
            "required": ["title", "slides"]
        }
    }
}


def _pptd_to_emu(ptpd_val: float, axis: str = "x") -> int:
    """将 PPTD 坐标/尺寸转为 EMU（English Metric Units）"""
    if axis == "x":
        return int(ptpd_val * Inches(PPTX_W_INCHES) / PPTD_W)
    else:
        return int(ptpd_val * Inches(PPTX_H_INCHES) / PPTD_H)


def _pptd_pos(ptpd_x: float, pptd_y: float) -> tuple:
    """将 PPTD 位置转为 EMU"""
    return (_pptd_to_emu(ptpd_x, "x"), _pptd_to_emu(pptd_y, "y"))


def _pptd_size(ptpd_w: float, pptd_h: float) -> tuple:
    """将 PPTD 尺寸转为 EMU"""
    return (_pptd_to_emu(ptpd_w, "x"), _pptd_to_emu(pptd_h, "y"))


def _parse_color(color_str: str) -> RGBColor:
    """解析颜色字符串，返回 RGBColor"""
    color_str = color_str.strip().lstrip("#")
    if len(color_str) == 3:
        color_str = "".join(c * 2 for c in color_str)
    r, g, b = int(color_str[0:2], 16), int(color_str[2:4], 16), int(color_str[4:6], 16)
    return RGBColor(r, g, b)


def _parse_style(style_str: str) -> dict:
    """解析 CSS style 字符串为字典"""
    if not style_str:
        return {}
    styles = {}
    for part in style_str.split(";"):
        part = part.strip()
        if ":" in part:
            key, val = part.split(":", 1)
            styles[key.strip()] = val.strip()
    return styles


def _apply_fill(shape, fill_value):
    """设置形状填充（支持纯色和渐变）"""
    if isinstance(fill_value, list):
        # 渐变填充
        try:
            fill = shape.fill
            fill.gradient()
            fill.gradient_angle = 90.0
            stops = fill.gradient_stops
            for i, stop_def in enumerate(fill_value):
                if isinstance(stop_def, dict):
                    color = _parse_color(stop_def.get("color", "#000000"))
                    offset = stop_def.get("offset", i / max(len(fill_value) - 1, 1))
                    if i < len(stops):
                        stops[i].color.rgb = color
                        stops[i].position = offset
        except Exception:
            _apply_solid_fill(shape, fill_value[0].get("color", "#000000") if fill_value else "#000000")
    elif isinstance(fill_value, str):
        _apply_solid_fill(shape, fill_value)


def _apply_solid_fill(shape, color_str: str):
    """设置纯色填充"""
    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _parse_color(color_str)
    except Exception:
        pass


def _parse_html_text(html_text: str) -> list:
    """
    解析 PPTD 的 HTML 文本，返回段落列表。
    每段格式: [(text, styles_dict), ...]
    """
    if not html_text:
        return []

    # 移除块级标签外的空白
    html_text = html_text.strip()

    paragraphs = []

    # 按 <p> 标签分割段落
    p_blocks = re.findall(r'<p\b[^>]*>(.*?)</p>', html_text, re.DOTALL)

    if not p_blocks:
        # 没有 <p> 标签，整个作为一段
        p_blocks = [html_text]

    for p_content in p_blocks:
        p_style = {}
        # 尝试从 <p> 标签中提取 style
        p_match = re.search(r'<p\b[^>]*style="([^"]*)"', html_text, re.DOTALL)
        if p_match:
            p_style = _parse_style(p_match.group(1))

        runs = []
        # 解析 <strong> 和 <span> 标签
        remaining = p_content.strip()
        tag_pattern = re.compile(r'<(strong|span)\b([^>]*)>(.*?)</\1>', re.DOTALL)

        last_end = 0
        for match in tag_pattern.finditer(remaining):
            # 前面的纯文本
            before = remaining[last_end:match.start()]
            if before.strip():
                runs.append((before, {}))
            last_end = match.end()

            tag = match.group(1)
            attrs = match.group(2)
            text = match.group(3)

            style = {}
            if tag == "strong":
                style["font-weight"] = "bold"
            style_attr = re.search(r'style="([^"]*)"', attrs)
            if style_attr:
                style.update(_parse_style(style_attr.group(1)))

            if text.strip():
                runs.append((text, style))

        # 最后剩余的纯文本
        after = remaining[last_end:]
        if after.strip():
            runs.append((after, {}))

        if runs:
            paragraphs.append({"runs": runs, "p_style": p_style})

    return paragraphs


def _add_text_element(slide, element: dict):
    """添加文本元素到幻灯片"""
    pos = element.get("position", [0, 0])
    size = element.get("size", [400, 100])
    content = element.get("content", {})
    html_text = content.get("text", "")

    left, top = _pptd_pos(pos[0], pos[1])
    width, height = _pptd_size(size[0], size[1])

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    paragraphs = _parse_html_text(html_text)

    for i, para in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p_style = para.get("p_style", {})

        # 段落级对齐
        align = p_style.get("text-align", "").lower()
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT

        for j, (text, run_style) in enumerate(para.get("runs", [])):
            if j == 0:
                run = p.runs[0] if p.runs else p.add_run()
                run.text = ""
            else:
                run = p.add_run()

            run.text = text

            # 合并段落样式和 run 样式
            merged_style = {**p_style, **run_style}

            # 字体大小
            font_size = merged_style.get("font-size", "")
            if font_size:
                try:
                    run.font.size = Pt(float(font_size.replace("px", "")))
                except ValueError:
                    pass

            # 字体颜色
            color = merged_style.get("color", "")
            if color and color.startswith("#"):
                try:
                    run.font.color.rgb = _parse_color(color)
                except Exception:
                    pass

            # 粗体
            if merged_style.get("font-weight") == "bold":
                run.font.bold = True

            # 斜体
            if merged_style.get("font-style") == "italic":
                run.font.italic = True

            # 下划线
            if merged_style.get("text-decoration") == "underline":
                run.font.underline = True

            # 字体
            font_family = merged_style.get("font-family", "")
            if font_family:
                run.font.name = font_family

            # 背景色
            bg = merged_style.get("background", "")
            if bg:
                try:
                    # 使用 XML 设置高亮
                    rPr = run._r.get_or_add_rPr()
                    highlight = rPr.makeelement(qn("a:highlight"), {})
                    srgb = highlight.makeelement(qn("a:srgbClr"), {"val": bg.lstrip("#")})
                    highlight.append(srgb)
                    rPr.append(highlight)
                except Exception:
                    pass

        # 段落间距（行高）
        line_height = p_style.get("line-height", "")
        if line_height:
            try:
                p.line_spacing = Pt(float(line_height) * 12)
            except ValueError:
                pass


def _add_shape_element(slide, element: dict):
    """添加形状元素到幻灯片"""
    pos = element.get("position", [0, 0])
    size = element.get("size", [100, 100])
    shape_def = element.get("shape", {})

    preset_name = shape_def.get("preset", "rectangle").lower().replace(" ", "")
    shape_type = SHAPE_MAP.get(preset_name, MSO_SHAPE.RECTANGLE)

    left, top = _pptd_pos(pos[0], pos[1])
    width, height = _pptd_size(size[0], size[1])

    shape = slide.shapes.add_shape(shape_type, left, top, width, height)

    # 填充
    fill_value = shape_def.get("fill", "")
    if fill_value:
        _apply_fill(shape, fill_value)

    # 边框
    stroke = shape_def.get("stroke", "")
    if stroke:
        shape.line.color.rgb = _parse_color(stroke)
        stroke_width = shape_def.get("strokeWidth", 1)
        shape.line.width = Pt(stroke_width)

    # 旋转
    rotation = shape_def.get("rotation", 0)
    if rotation:
        shape.rotation = rotation

    # 文本（形状内嵌文字）
    text = shape_def.get("text", "")
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(12)


def _add_image_element(slide, element: dict):
    """添加图片元素到幻灯片"""
    pos = element.get("position", [0, 0])
    size = element.get("size", [200, 150])
    image_def = element.get("image", {})
    src = image_def.get("src", "")

    if not src:
        return

    left, top = _pptd_pos(pos[0], pos[1])
    width, height = _pptd_size(size[0], size[1])

    try:
        if src.startswith(("http://", "https://")):
            import urllib.request
            data = urllib.request.urlopen(src, timeout=30).read()
            slide.shapes.add_picture(BytesIO(data), left, top, width, height)
        elif os.path.exists(src):
            slide.shapes.add_picture(src, left, top, width, height)
    except Exception as e:
        # 图片加载失败，添加占位符
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        tf = shape.text_frame
        tf.paragraphs[0].text = f"[图片: {src}]"
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = RGBColor(0x99, 0x99, 0x99)


def _render_slide(slide, elements: list, theme_colors: list):
    """渲染一页幻灯片的所有元素"""
    for element in elements:
        if not isinstance(element, dict):
            continue
        elem_type = element.get("type", "").lower()

        if elem_type == "text":
            _add_text_element(slide, element)
        elif elem_type == "shape":
            _add_shape_element(slide, element)
        elif elem_type == "image":
            _add_image_element(slide, element)


def _parse_page_content(content: str) -> list:
    """解析 .page 文件的 YAML 内容，返回 elements 列表"""
    if not content or not content.strip():
        return []
    try:
        data = yaml.safe_load(content)
        if isinstance(data, dict) and "elements" in data:
            return data["elements"]
        return []
    except yaml.YAMLError:
        return []


def generate_kimi_ppt(title: str, theme_colors: list = None, slides: list = None) -> dict:
    """
    生成 PPT，使用 python-pptx 本地渲染 PPTD 格式

    Args:
        title: PPT 标题
        theme_colors: 主题配色（5个颜色）
        slides: 页面列表，每项包含 filename 和 content（YAML 字符串）

    Returns:
        dict: {"success": bool, "url": str, "filepath": str, "message": str}
    """
    if not slides:
        return {"success": False, "error": "slides 不能为空"}

    if theme_colors is None:
        theme_colors = ["#2B579A", "#E86A17", "#333333", "#FFFFFF", "#F0F2F5"]

    project_id = uuid.uuid4().hex[:8]
    project_name = f"kimi_ppt_{project_id}"

    try:
        emit_progress("generate_kimi_ppt", 10, "正在初始化演示文稿...")
        # 创建 PPTX
        prs = Presentation()
        prs.slide_width = Inches(PPTX_W_INCHES)
        prs.slide_height = Inches(PPTX_H_INCHES)

        # 使用空白布局
        blank_layout = prs.slide_layouts[6]  # blank

        # 渲染每一页（按页上报进度）
        total = len(slides)
        for i, slide_data in enumerate(slides):
            pct = 10 + int((i + 1) / max(total, 1) * 80)
            emit_progress("generate_kimi_ppt", pct, f"正在渲染第 {i + 1}/{total} 页...")
            slide = prs.slides.add_slide(blank_layout)

            # 设置幻灯片背景
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = _parse_color(theme_colors[4])  # 浅灰背景

            content = slide_data.get("content", "")
            elements = _parse_page_content(content)
            _render_slide(slide, elements, theme_colors)

        # 保存 PPTX
        emit_progress("generate_kimi_ppt", 95, "正在保存文件...")
        output_filename = f"{project_name}.pptx"
        output_path = OUTPUT_DIR / output_filename
        prs.save(str(output_path))
        emit_progress("generate_kimi_ppt", 100, "PPT 生成完成")

        url = f"/static/papers/{output_filename}"

        return {
            "success": True,
            "url": url,
            "filepath": str(output_path),
            "filename": output_filename,
            "message": (
                f"PPT 已生成（PPTD 本地渲染），共 {len(slides)} 页，"
                f"下载链接：[PAPER:{url}]"
            )
        }

    except Exception as e:
        return {"success": False, "error": str(e)}


def execute(arguments: dict) -> str:
    """工具统一入口，由 tools/__init__.py 调度"""
    result = generate_kimi_ppt(
        title=arguments.get("title", ""),
        theme_colors=arguments.get("theme_colors"),
        slides=arguments.get("slides", [])
    )
    if result.get("success"):
        return result["message"]
    else:
        return f"PPT 生成失败：{result.get('error', '未知错误')}"