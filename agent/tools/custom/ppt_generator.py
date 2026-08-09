"""
PPT 生成工具 —— 使用 python-pptx 生成演示文稿
输出文件存储：server/static/papers/
"""

import os
import uuid

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from agent.tools import emit_progress

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))
OUTPUT_DIR = os.path.join(ROOT, "server", "static", "papers")
os.makedirs(OUTPUT_DIR, exist_ok=True)

# 配色方案
COLOR_PRIMARY = RGBColor(0x2B, 0x57, 0x9A)   # 深蓝
COLOR_ACCENT = RGBColor(0xE8, 0x6A, 0x17)    # 橙色
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DARK = RGBColor(0x33, 0x33, 0x33)
COLOR_LIGHT_BG = RGBColor(0xF0, 0xF2, 0xF5)

SCHEMA = {
    "type": "function",
    "tag": "PPT",
    "function": {
        "name": "generate_ppt",
        "description": "生成 PPT 演示文稿。title 为标题，slides 为每页内容要点数组。",
        "parameters": {
            "type": "object",
            "properties": {
                "title": {
                    "type": "string",
                    "description": "PPT的标题，显示在封面页"
                },
                "subtitle": {
                    "type": "string",
                    "description": "PPT的副标题，显示在封面页标题下方"
                },
                "slides": {
                    "type": "array",
                    "description": "幻灯片内容列表，每个元素包含一页的标题和内容要点",
                    "items": {
                        "type": "object",
                        "properties": {
                            "slide_title": {
                                "type": "string",
                                "description": "该页幻灯片的标题"
                            },
                            "bullets": {
                                "type": "array",
                                "description": "该页的内容要点列表，每个元素为一条要点",
                                "items": {"type": "string"}
                            }
                        },
                        "required": ["slide_title", "bullets"]
                    }
                }
            },
            "required": ["title", "slides"]
        }
    }
}


def _add_cover_slide(prs, title, subtitle=""):
    """添加封面页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # 顶部色块
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0), Inches(0), slide_w, Inches(0.15)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()

    # 左侧色带
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(0.08), slide_h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.fill.background()

    # 标题
    left = Inches(1.2)
    top = Inches(2.0)
    width = slide_w - Inches(2.4)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.alignment = PP_ALIGN.LEFT

    # 副标题
    if subtitle:
        top2 = Inches(3.5)
        txBox2 = slide.shapes.add_textbox(left, top2, width, Inches(1.0))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = COLOR_DARK
        p2.alignment = PP_ALIGN.LEFT

    # 底部装饰线
    shape = slide.shapes.add_shape(
        1, Inches(0), slide_h - Inches(0.05), slide_w, Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.fill.background()


def _add_content_slide(prs, slide_title, bullets):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # 顶部色条
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), slide_w, Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()

    # 标题背景
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0.08), slide_w, Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()

    # 标题文字
    txBox = slide.shapes.add_textbox(
        Inches(1.0), Inches(0.2), slide_w - Inches(2.0), Inches(0.9)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.LEFT

    # 内容区
    content_top = Inches(1.6)
    content_left = Inches(1.2)
    content_width = slide_w - Inches(2.4)
    content_height = slide_h - Inches(2.2)

    txBox = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_DARK
        p.space_after = Pt(12)
        p.alignment = PP_ALIGN.LEFT

    # 底部页码装饰
    shape = slide.shapes.add_shape(
        1, Inches(0), slide_h - Inches(0.05), slide_w, Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.fill.background()


def _add_end_slide(prs):
    """添加结束页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # 背景
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), slide_w, slide_h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()

    # 感谢文字
    txBox = slide.shapes.add_textbox(
        Inches(0), Inches(2.5), slide_w, Inches(1.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "感谢聆听"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # 装饰线
    shape = slide.shapes.add_shape(
        1, Inches(3.5), Inches(4.2), Inches(3.0), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.fill.background()


def generate_ppt(title: str, subtitle: str = "", slides: list = None) -> dict:
    """
    生成PPT演示文稿

    Args:
        title: PPT标题
        subtitle: 副标题（可选）
        slides: 幻灯片列表，每项包含 slide_title 和 bullets

    Returns:
        dict: {"success": bool, "url": str, "filepath": str}
    """
    if not slides:
        return {"success": False, "error": "slides 不能为空"}

    try:
        emit_progress("generate_ppt", 5, "正在初始化演示文稿...")
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # 封面
        _add_cover_slide(prs, title, subtitle)

        # 内容页（按页上报进度）
        total = len(slides)
        for i, slide_data in enumerate(slides):
            pct = 10 + int((i + 1) / max(total, 1) * 75)
            emit_progress("generate_ppt", pct, f"正在排版第 {i + 1}/{total} 页...")
            _add_content_slide(
                prs,
                slide_data.get("slide_title", ""),
                slide_data.get("bullets", [])
            )

        # 结束页
        _add_end_slide(prs)

        # 保存
        emit_progress("generate_ppt", 90, "正在保存文件...")
        filename = f"ppt_{uuid.uuid4().hex[:8]}.pptx"
        filepath = os.path.join(OUTPUT_DIR, filename)
        prs.save(filepath)
        emit_progress("generate_ppt", 100, "PPT 生成完成")

        url = f"/static/papers/{filename}"

        return {
            "success": True,
            "url": url,
            "filepath": filepath,
            "filename": filename,
            "message": f"PPT已生成，共 {len(slides) + 2} 页（含封面和结束页），下载链接：[PAPER:{url}]"
        }
    except Exception as e:
        return {"success": False, "error": str(e)}


def execute(arguments: dict) -> str:
    """工具统一入口，由 tools/__init__.py 调度"""
    import json
    result = generate_ppt(
        title=arguments.get("title", ""),
        subtitle=arguments.get("subtitle", ""),
        slides=arguments.get("slides", [])
    )
    return json.dumps(result, ensure_ascii=False)